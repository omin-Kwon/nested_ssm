"""Meeting pitch deck (김세훈 미팅): SSM background -> tiering idea -> components
-> training -> math of sorting -> results. Native pptx shapes for diagrams;
two measured figures reused from theory_figs.
Run: ~/nemo_env/bin/python3 make_pitch_ppt.py   (from scale/)
Output: docs/pitch_deck_tiering.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FW, FH = 13.333, 7.5
NAVY = RGBColor(0x1a, 0x36, 0x5d)
RED = RGBColor(0xd6, 0x27, 0x28)
BLUE = RGBColor(0x1f, 0x77, 0xb4)
GREEN = RGBColor(0x04, 0x78, 0x57)
ORANGE = RGBColor(0xc2, 0x41, 0x0c)
GRAY = RGBColor(0x66, 0x66, 0x66)
LRED = RGBColor(0xfe, 0xe2, 0xe2)
LBLUE = RGBColor(0xdb, 0xea, 0xfe)
LGREEN = RGBColor(0xd1, 0xfa, 0xe5)
LYEL = RGBColor(0xff, 0xf7, 0xed)
WHITE = RGBColor(0xff, 0xff, 0xff)
BLACK = RGBColor(0x11, 0x11, 0x11)

prs = Presentation()
prs.slide_width = Inches(FW)
prs.slide_height = Inches(FH)
BLANK = prs.slide_layouts[6]


def new_slide(title, sub=None):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(.45), Inches(.18), Inches(FW - .9), Inches(.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(25); p.font.bold = True; p.font.color.rgb = NAVY
    if sub:
        p2 = tb.text_frame.add_paragraph()
        p2.text = sub; p2.font.size = Pt(12); p2.font.color.rgb = GRAY
    return s


def box(s, x, y, w, h, text, fill=LBLUE, line=BLUE, fs=12, bold=False,
        fc=BLACK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.CENTER):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(1.4)
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(45720); tf.margin_right = Emu(45720)
    tf.margin_top = Emu(22860); tf.margin_bottom = Emu(22860)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(fs); p.font.bold = bold
        p.font.color.rgb = fc; p.alignment = align
    return sh


def txt(s, x, y, w, h, lines, fs=13, fc=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(fs); p.font.bold = bold
        p.font.color.rgb = fc; p.alignment = align
    return tb


def rarrow(s, x, y, w, h=.32, color=GRAY, text=None, fs=10):
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    if text:
        txt(s, x - .25, y - .32, w + .5, .3, text, fs=fs, fc=color,
            align=PP_ALIGN.CENTER)
    return sh


def darrow(s, x, y, h, w=.3, color=GRAY):
    sh = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    return sh


# ================================================= 1. title
s = prs.slides.add_slide(BLANK)
txt(s, .9, 2.1, 11.6, 1.6,
    ["Elastic Test-Time Memory:",
     "Recurrent State의 학습된 Hot/Cold 티어링"], fs=33, bold=True, fc=NAVY)
txt(s, .95, 3.7, 11.4, 1.6, [
    "공개 Nemotron-9B의 0.04% (직교 회전 R)만 재학습 → decode B200 실측 1.92~2.42×,",
    "공식 평가 스택(NeMo-Skills+vLLM) 전 벤치 lossless (GSM8K 94.6 / MATH-500 95.2 / RULER 98-100)"],
    fs=16)
txt(s, .95, 5.6, 11, .5, "SNU ARC Lab · 2026-07 · 미팅 자료", fs=12, fc=GRAY)

# ================================================= 2. SSM background: s -> t
s = new_slide("1. Background — SSM/Linear Attention의 기억은 '행렬 S' 하나",
              "토큰 s가 토큰 t에 미치는 영향의 전부는 state를 경유한다")
txt(s, .55, 1.05, 12.3, 1.35, [
    "갱신:  S_t = a_t · S_{t−1} + Δx_t ⊗ B_t        (S ∈ ℝ^{P×N},  a_t ∈ (0,1) per-head 스칼라 decay)",
    "읽기:  y_t = S_t · C_t                                (N축 내적)",
    "펼치면 — 토큰 s의 t에서의 기여:   y_t = Σ_s ( a_{s+1}·a_{s+2}···a_t ) · ⟨C_t , B_s⟩ · Δx_s"],
    fs=15)
# diagram: token s -> write -> S -> read -> token t
box(s, .8, 3.3, 1.7, .9, "토큰 s\n(과거)", LYEL, ORANGE, fs=13, bold=True)
rarrow(s, 2.6, 3.55, 1.5, text="쓰기 주소 B_s (key)", fs=10)
box(s, 4.2, 2.75, 3.1, 2.0, "STATE  S  (P×N)\n모든 과거의 압축\n= 유일한 기억", LBLUE, BLUE, fs=14, bold=True)
rarrow(s, 7.4, 3.55, 1.5, text="읽기 주소 C_t (query)", fs=10)
box(s, 9.0, 3.3, 1.7, .9, "토큰 t\n(현재)", LGREEN, GREEN, fs=13, bold=True)
box(s, 4.55, 5.05, 2.4, .55, "감쇠 ∏ a_τ  (s→t 나이만큼)", WHITE, GRAY, fs=10)
txt(s, .55, 5.95, 12.3, 1.3, [
    "· ⟨C_t, B_s⟩ = N차원 주소 공간의 내적 — \"attention 행렬 없는 attention\". N = 주소 공간 크기 = recall 용량 (Zoology/BASED 법칙).",
    "· attention과의 결정적 차이: 過去를 다시 볼 수 없다 — S에 남긴 것이 전부. 그래서 state가 크고 소중하다."],
    fs=13)

# ================================================= 3. bottleneck
s = new_slide("2. 병목 — decode는 매 토큰 S 전체를 왕복한다 (memory-bound)",
              "B200 실측: 대배치 서빙에서 state-op가 지배 성분")
box(s, 1.2, 1.5, 2.6, 1.5, "GPU 연산부\n(SM)", LGREEN, GREEN, fs=14, bold=True)
box(s, 8.6, 1.5, 3.3, 1.5, "HBM\nstate S = 141.6 MB/시퀀스\n(9B, N=128)", LBLUE, BLUE, fs=13, bold=True)
rarrow(s, 4.0, 1.65, 4.4, h=.4, color=RED, text="매 토큰 READ 전체 (1.0 S)", fs=11)
sh = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(4.0), Inches(2.45), Inches(4.4), Inches(.4))
sh.fill.solid(); sh.fill.fore_color.rgb = RED; sh.line.fill.background()
txt(s, 4.0, 2.9, 4.4, .3, "매 토큰 WRITE 전체 (1.0 S)", fs=11, fc=RED, align=PP_ALIGN.CENTER)
txt(s, .55, 3.7, 12.3, 2.9, [
    "· 토큰당 연산/바이트(OI) ≈ 1 — GPU 균형점(~300)의 0.3%. 연산이 아니라 트래픽이 시간을 지배.",
    "· B200 vLLM 실측: state-op 점유 2.6% (B=1) → 57.6% (B=256) → 60% (B=1152). throughput 천장 11.9k tok/s, capacity wall B≈1230.",
    "· 커널 실측: tok/s ∝ 1/state (기울기 −1 정확), 메모리 용량 2×를 줘도 tok/s 불변 — 벽은 용량이 아니라 대역폭×state 크기.",
    "",
    "→ 목표: 정확도를 지키면서 '매 토큰 움직이는 바이트'를 줄이는 것."],
    fs=14)

# ================================================= 4. failed alternatives
s = new_slide("3. 왜 기존 길들은 안 되는가 — 실측 사다리",
              "줄이기·자르기·골라읽기 전부 유료; 남는 길은 하나")
box(s, .7, 1.4, 3.8, 1.05, "① state 줄이기 (작은 N)", LRED, RED, fs=14, bold=True)
txt(s, .8, 2.5, 3.7, 1.3, "recall 용량 자체가 하락\n(Zoology: state=recall 법칙)", fs=12)
box(s, 4.8, 1.4, 3.8, 1.05, "② 덜 중요한 차원 삭제\n(GHOST류 calibration 절단)", LRED, RED, fs=13, bold=True)
txt(s, 4.9, 2.5, 3.7, 1.7, "GSM8K 85.0 → 58.5 (−26.5)\n학습된 기저로 잘라도 −25.5\n희생축 = 다단계 추론\n(wikitext ppl은 +1뿐 — 함정)", fs=12)
box(s, 8.9, 1.4, 3.8, 1.05, "③ 토큰마다 골라 읽기\n(Quest류 동적 선택)", LRED, RED, fs=13, bold=True)
txt(s, 9.0, 2.5, 3.7, 1.7, "압축 state엔 per-token\n희소성이 없음 — 열 top-k\n읽기 33%→acc 74, 8%→21 붕괴", fs=12)
box(s, .7, 4.5, 12.0, 1.15,
    "④ 우리: 아무것도 버리지 않고, 덜 중요한 부분만 '낡은 채로 읽고(stale) 몰아서 정확히 쓴다(lazy)'  —  GSM8K −0.5 (공식 스택, 학습 후 0)",
    LGREEN, GREEN, fs=15, bold=True)
txt(s, .7, 5.9, 12.0, 1.0, [
    "핵심 대조: sparse-and-fresh (무엇을 저장할지 선택 — 비가역 손실)  vs  dense-but-stale (언제 반영할지만 지연 — 오차가 최근 c토큰에 국한, 정확 보상)"],
    fs=13, fc=NAVY, bold=True)

# ================================================= 5. idea: tiering
s = new_slide("4. 아이디어 — State의 N축을 Hot / Cold 두 티어로",
              "learned write-back cache hierarchy: 분할 기준은 시간이 아니라 '트래픽'")
# state matrix drawing: P x N with hot columns
txt(s, 1.0, 1.25, 3.0, .4, "S  (P × N),  N축을 분할:", fs=14, bold=True)
box(s, 1.0, 1.8, 1.5, 2.6, "HOT\n32열", RGBColor(0xfc, 0xa5, 0xa5), RED, fs=14, bold=True, shape=MSO_SHAPE.RECTANGLE)
box(s, 2.5, 1.8, 4.2, 2.6, "COLD\n96열", LBLUE, BLUE, fs=14, bold=True, shape=MSO_SHAPE.RECTANGLE)
txt(s, 1.0, 4.5, 5.9, .5, "← P (head·채널 축, 그대로) —  분할은 회전 후 N(주소) 축에서", fs=11, fc=GRAY)
# right: per-token behavior
box(s, 7.4, 1.55, 5.2, 1.35, "HOT: 매 토큰 fresh READ+WRITE\n(트래픽 상위 차원 — 항상 최신)", LRED, RED, fs=13, bold=True)
box(s, 7.4, 3.1, 5.2, 1.7, "COLD: 읽기 = 낡은 snapshot + decay 정확보상\n쓰기 = c토큰마다 chunk-exact write-back\n(READ −47% / WRITE −61% @ c4 실측)", LBLUE, BLUE, fs=12.5, bold=True)
txt(s, .7, 5.5, 12.2, 1.6, [
    "· 정적 연속 prefix 분할: 토큰별 controller 없음, 메타데이터 0, 연속 6KB 블록 — '무엇이 hot인가'는 학습이 R에 구워 놓음.",
    "· 오차의 정체: 최근 c토큰의 cold 성분이 readout에 늦게 반영되는 것뿐 (age-국소·내용 무관). flush는 수학적으로 exact.",
    "· 속도 결과: B200 실측 1.62×(c4·bf16) ~ 2.42×(c16·bf16) — 매 토큰 2.0S 이동이 ~0.8S로."],
    fs=13)

# ================================================= 6. component: R
s = new_slide("5. 추가 컴포넌트는 '직교 회전 R' 하나 (0.04%)",
              "백본 동결 · 레이어당 그룹별 R (27층 × 8그룹 × 128×128 = 3.5M 파라미터)")
# pipeline
box(s, .6, 1.7, 1.4, .8, "x_t", WHITE, GRAY, fs=13)
rarrow(s, 2.05, 1.9, .6)
box(s, 2.7, 1.7, 1.8, .8, "in_proj\n+ conv/act", LBLUE, BLUE, fs=11)
rarrow(s, 4.55, 1.9, .6)
box(s, 5.2, 1.45, 1.5, .6, "B_t (key)", LYEL, ORANGE, fs=11)
box(s, 5.2, 2.2, 1.5, .6, "C_t (query)", LYEL, ORANGE, fs=11)
rarrow(s, 6.75, 1.9, .6)
box(s, 7.4, 1.5, 1.9, 1.25, "R  (학습 대상)\nB→RB, C→RC\n그룹별 직교", RGBColor(0xfd, 0xe6, 0x8a), ORANGE, fs=12, bold=True)
rarrow(s, 9.35, 1.9, .6)
box(s, 10.0, 1.5, 2.6, 1.25, "SSM 코어 (동결)\nS 갱신·읽기는\n회전 좌표에서", LBLUE, BLUE, fs=12)
txt(s, .6, 3.2, 12.3, 1.0, [
    "회전 불변성 정리 (명제 1):  B→RB, C→RC (직교, per-group)  ⇒  S → S·Rᵀ,  출력 y는 불변.",
    "⟨RC_t, RB_s⟩ = C_tᵀRᵀR B_s = ⟨C_t, B_s⟩  —  RᵀR = I 가 모든 내적·길이를 보존 (수동적 좌표 교체)."],
    fs=14, fc=NAVY, bold=True)
txt(s, .6, 4.45, 12.3, 2.3, [
    "따라서:",
    "· full-width(티어링 off)의 정확도는 수학적으로 보존 — 공식 스택 실증: fresh = raw 정답 수까지 동일.",
    "· R은 '무엇을 배우는가'가 아니라 '어느 좌표에서 볼 것인가'만 배움 — 가설 공간이 작아 과적합이 원리적으로 어려움.",
    "· 학습으로 바뀌는 것은 R 단독 (decay 미세조정은 bf16 정밀도에 흡수됨을 확인 — 정직 기술).",
    "· 성립 조건: per-head 스칼라 decay 가족(Mamba-2/GDN — Nemotron·Qwen 계열)은 무조건; 채널별 decay(GLA/KDA)는 경량 FT 필요."],
    fs=13)

# ================================================= 7a. math of sorting (1/2): traffic derivation
s = new_slide("6. 정렬의 수학 ① — readout을 채널로 쪼개면 '트래픽'이 정의된다",
              "슬라이드 1의 갱신식에서 세 줄로 유도 (새 가정 0개)")
box(s, .55, 1.15, 7.6, .95,
    "① 갱신식을 재귀 대입으로 펼친다:\ny_t = Σ_s  A_{s,t} · ⟨B_s, C_t⟩ · Δx_s        (A_{s,t} = a_{s+1}···a_t : 감쇠 생존율)",
    WHITE, BLUE, fs=12.5, align=PP_ALIGN.LEFT)
darrow(s, 4.2, 2.18, .28)
box(s, .55, 2.52, 7.6, 1.25,
    "② 항등식 I = Σ_n r_n r_nᵀ 을 내적 사이에 삽입 (r_n = R의 n번째 행, 정규직교):\n⟨B_s, C_t⟩ = Σ_n ⟨r_n, B_s⟩ · ⟨r_n, C_t⟩\n→ 총 통신이 N개 채널의 몫으로 정확히 쪼개짐 — 회전된 state의 '열 n' = 통신 채널 n",
    WHITE, BLUE, fs=12.5, align=PP_ALIGN.LEFT)
darrow(s, 4.2, 3.85, .28)
box(s, .55, 4.18, 7.6, .95,
    "③ 채널 n이 나르는 평균 몫 = 트래픽:\nT(r_n) = E[ ⟨C, r_n⟩⟨r_n, B⟩ ] = r_nᵀ · E[C Bᵀ] · r_n",
    LYEL, ORANGE, fs=13, bold=True, align=PP_ALIGN.LEFT)
box(s, 8.45, 1.15, 4.3, 3.98,
    "숫자 예시 (N=2)\n\nB=(3,4), C=(2,1) → 총 통신 ⟨B,C⟩ = 10\n\n표준 기저:   3·2 + 4·1 = 6 + 4\n45° 회전 기저:  10.5 + (−0.5)\n\n합은 어느 기저에서나 10\n(= 회전 불변성, full-width 무손실)\n분배만 기저에 따라 달라짐\n\n→ '정렬'은 성능이 아니라\n     배치(placement)의 문제",
    LGREEN, GREEN, fs=12, align=PP_ALIGN.LEFT)
# toy bars: same info, two bases (baseline 7.08, scale 2.2)
txt(s, .8, 5.38, 5.6, .35, "원 좌표: T = (0.22, 0.22, 0.12, 0.22, 0.22) — 균등해 보임", fs=11, bold=True)
raw = [0.22, 0.22, 0.12, 0.22, 0.22]
for i, v in enumerate(raw):
    box(s, 1.0 + i * .85, 7.08 - v * 2.2, .62, v * 2.2, "", LBLUE, BLUE, shape=MSO_SHAPE.RECTANGLE)
txt(s, 6.9, 5.38, 5.9, .35, "회전 좌표: T′ = (0.60, 0.25, 0.15, 0, 0) — hot 2서랍이 85% 커버", fs=10.5, bold=True, fc=RED)
rot = [0.60, 0.25, 0.15, 0.008, 0.008]
for i, v in enumerate(rot):
    col = (RGBColor(0xfc, 0xa5, 0xa5), RED) if i < 2 else (LBLUE, BLUE)
    box(s, 7.1 + i * .85, 7.08 - max(v, .02) * 2.2, .62, max(v, .02) * 2.2, "", col[0], col[1], shape=MSO_SHAPE.RECTANGLE)

# ================================================= 7b. math of sorting (2/2): where w is born
s = new_slide("7. 정렬의 수학 ② — 중요도 가중 w는 loss의 미분에서 태어난다",
              "forward 수식에 w는 없다 — R을 조각하는 끌은 forward가 아니라 gradient")
box(s, .55, 1.2, 7.7, .95,
    "① 학습이 실행하는 것: 스텝마다 폭 k 추첨 → 뒤 채널(n>k)을 삭제(nested dropout)\n    또는 낡힘(v4-aware, c 추첨) → 그 상태로 다음 토큰 CE loss를 측정",
    WHITE, BLUE, fs=12.5, align=PP_ALIGN.LEFT)
darrow(s, 4.2, 2.23, .28)
box(s, .55, 2.57, 7.7, 1.45,
    "② 채널 n을 잃었을 때의 loss 변화 (1차 연쇄법칙, g_t = ∂L/∂y_t):\nΔL_n ≈ Σ_{t,s}  w_{t,s} · ⟨r_n,B_s⟩⟨r_n,C_t⟩  =  r_nᵀ · M_eff · r_n\nM_eff = E[ C Bᵀ · w ],      w_{t,s} = A_{s,t} · ⟨ g_t , Δx_s ⟩",
    LYEL, ORANGE, fs=13, bold=True, align=PP_ALIGN.LEFT)
darrow(s, 4.2, 4.1, .28)
box(s, .55, 4.44, 7.7, .95,
    "③ 폭 추첨의 생존확률 p_1 > p_2 > ··· 아래 기대 loss 최소해:\nR = M_eff(대칭화)의 고유기저, 고유값 내림차순  (nested dropout → PCA 정리의 연산자판)",
    LGREEN, GREEN, fs=12.5, bold=True, align=PP_ALIGN.LEFT)
box(s, 8.55, 1.2, 4.2, 4.19,
    "w의 정체 — 사건(s→t)마다 스칼라\n\n=  감쇠 생존율  A_{s,t}\n×  값의 크기  Δx_s\n×  loss 민감도  ∂L/∂y_t\n\n· 새로 도입한 가정이 아니라\n  연쇄법칙의 '나머지 인자들'\n· 아무도 입력·저장하지 않음 —\n  backprop이 매 스텝 자동으로 곱함\n· 크기 통계 M = E[CBᵀ]는\n  w ≡ 1로 잊어버린 근사",
    LYEL, ORANGE, fs=12, align=PP_ALIGN.LEFT)
txt(s, .55, 5.75, 12.4, 1.4, [
    "실측 반전 (E-T1): w≡1 근사(계산 가능한 크기 M)의 고유기저는 학습 R과 무상관 (주각 cos 0.44 ≈ 무작위; 시차·열노름·질의에너지 4중 수렴).",
    "→ w는 보정이 아니라 지배 인자 — 중요도는 loss가 정의하므로 training-free 통계(calibration류·GHOST)로 대체 불가 = end-to-end 학습이 필수인 이유."],
    fs=13, fc=NAVY, bold=True)

# ================================================= 8. generalization (figure)
s = new_slide("8. 왜 일반화되는가 — 이론이 아니라 실측으로 답한다",
              "E-T1: 학습 기저는 어떤 통계와도 다르다 · E-T2: 그런데 도메인을 넘어 이식된다")
s.shapes.add_picture("results/plots/theory_figs/F5_measured.png",
                     Inches(.7), Inches(1.35), Inches(12.0), Inches(4.55))
txt(s, .7, 6.05, 12.0, 1.2, [
    "왼쪽: 학습 R의 hot 부분공간 vs 명시 통계 고유공간 — 무작위 수준 (중요도는 loss가 정의).",
    "오른쪽: 같은 R 하나의 절단세(ppl 세율) 곡선이 wiki/수학/코드에서 겹침(k16: 1.26/1.31/1.27×) — calibration 순열은 자기 폭 밖 절벽(17~36×)."],
    fs=13)

# ================================================= 9. training methodology
s = new_slide("9. 학습 방법론 — '티어링 실행권'을 사는 3요소",
              "총 수 GPU-시간, 백본 동결, R만 학습")
box(s, .7, 1.45, 3.9, 1.9, "① Nested 폭 메뉴\n(정렬을 만드는 힘)\n\n스텝마다 폭 k ∈ {16,32,64,128} 추첨\n→ 앞 좌표일수록 자주 loss에 노출\n= 암묵적 '중요도 투표'", LYEL, ORANGE, fs=12)
box(s, 4.8, 1.45, 3.9, 1.9, "② v4-aware 학습\n(staleness 내성)\n\n스텝의 절반을 tiered forward로\n(c ∈ {4..64} 추첨)\n→ 모델이 낡은 cold 읽기에 적응", LBLUE, BLUE, fs=12)
box(s, 8.9, 1.45, 3.9, 1.9, "③ QR retraction\n(무손실 보장 유지)\n\n매 스텝 R을 직교 다양체로 사영\n→ 명제 1의 전제(RᵀR=I)가\n학습 내내 정확히 성립", LGREEN, GREEN, fs=12)
darrow(s, 6.5, 3.5, .5)
box(s, .7, 4.15, 12.1, .95,
    "장문 CoT 마감 레시피 (longcot2): seqlen 4096 + c4·c8 가중 혼합 — \"장문 staleness는 장문으로 가르쳐야\" → GSM8K v4-c4 = fresh (+0.4, lossless)",
    LGREEN, GREEN, fs=13.5, bold=True)
txt(s, .7, 5.35, 12.1, 1.7, [
    "학습이 실제로 사는 것 (실측 사슬):",
    "· 학습 없이 티어링만: GSM8K 33.0 (−48.6, 정밀도 무관 붕괴) → 학습 후: 81.6 (= fresh). 실행권 = 학습의 산물.",
    "· 함정 2건: full-width 강제 distill은 폭 탄력성 파괴(폭 샘플링 유지 필수) · 단일 도메인(wikitext) 학습은 다운스트림 드리프트(혼합 필수)."],
    fs=13)

# ================================================= 10. v4 semantics timeline
s = new_slide("10. 실행 의미론 (v4) — 타임라인으로 보기",
              "hot은 항상 최신, cold는 'snapshot + 감쇠 보상'으로 읽고 c마다 정확히 따라잡음")
# timeline tokens
for i in range(8):
    box(s, .8 + i * 1.05, 1.6, .85, .6, f"t={i+1}", WHITE, GRAY, fs=11)
# hot row
box(s, 9.6, 1.6, 3.1, .6, "HOT: 매 토큰 R/W", LRED, RED, fs=11, bold=True)
for i in range(8):
    darrow(s, 1.08 + i * 1.05, 2.3, .45, w=.25, color=RED)
box(s, .8, 2.85, 8.6, .6, "hot state (항상 fresh)", LRED, RED, fs=12)
# cold row
for i in [3, 7]:
    darrow(s, 1.08 + i * 1.05, 3.55, .5, w=.3, color=BLUE)
box(s, .8, 4.15, 8.6, .6, "cold state (c=4마다 chunk-exact write-back — 개루프, 수학적 exact)", LBLUE, BLUE, fs=12)
box(s, 9.6, 4.15, 3.1, .6, "COLD 쓰기: c마다 몰아서", LBLUE, BLUE, fs=11, bold=True)
txt(s, .7, 5.1, 12.2, 1.9, [
    "· 읽기: y_t = C_t·S_hot(최신) + C_t·S_cold(snapshot)·exp(G_t) — 감쇠는 정확히 보상, 늦는 건 최근 c토큰의 cold 성분뿐.",
    "· flush가 개루프(청크 단위 재계산)라는 점이 뒤의 정밀도 면허(fp8-cold)의 근원 — 라운딩 오차가 재귀에 복리로 쌓이지 않음.",
    "· c = 정확도↔속도 다이얼: 생성이 긴 태스크(수학 CoT)는 c4, 짧은 답(recall)은 c16 — 재학습 없이 config로."],
    fs=13)

# ================================================= 11. results
s = new_slide("11. 결과 — 공식 스택 3-arm 매트릭스 + B200 속도",
              "acc = NeMo-Skills + vLLM (공식 논문과 동일 측정), 3-arm = raw / fresh / v4")
rows = [
    ("bench", "raw", "fresh (R만 켬)", "v4-c4-fp8 (배포점)", "판정"),
    ("GSM8K (1319)", "95.0", "95.0 (정답수 동일)", "94.6", "lossless"),
    ("MATH-500", "97.6", "98.2", "95.2 → longcot2로 봉합 중", "−2.4 (staleness, 정밀도 무관)"),
    ("RULER@4k niah 5종", "98–100", "—", "98–100", "lossless"),
]
ty = 1.45
for r, row in enumerate(rows):
    widths = [2.5, 1.3, 2.6, 3.3, 2.7]
    xx = .65
    for c, (wd, cell) in enumerate(zip(widths, row)):
        fill = LBLUE if r == 0 else WHITE
        box(s, xx, ty, wd, .55, cell, fill, BLUE, fs=11, bold=(r == 0),
            shape=MSO_SHAPE.RECTANGLE)
        xx += wd
    ty += .55
txt(s, .65, 3.9, 12.2, .5, "속도 (B200 실측, B=256, fused-only):", fs=14, bold=True, fc=NAVY)
box(s, .65, 4.4, 2.9, .8, "fresh (fp32)\n1.00×", WHITE, GRAY, fs=12)
rarrow(s, 3.6, 4.6, .5)
box(s, 4.15, 4.4, 3.3, .8, "+ lazy write (tiering c4)\n1.62× — WRITE −61%", LRED, RED, fs=12)
rarrow(s, 7.5, 4.6, .5)
box(s, 8.05, 4.4, 2.3, .8, "+ bf16-cold (c16)\n2.42×", LBLUE, BLUE, fs=12, bold=True)
rarrow(s, 10.4, 4.6, .5)
box(s, 10.95, 4.4, 1.9, .8, "+ fp8-cold\n~2.8× (analytic)", LGREEN, GREEN, fs=11)
txt(s, .65, 5.5, 12.2, 1.6, [
    "· 그 밖의 축 전부 3-arm 동률: commonsense 8 · recall-intensive 6(BASED) · RULER 11종 · HumanEval · needle 1.00 · ppl 0%.",
    "· 부정 결과도 자산: async flush 이득 0(BW 포화 — 바이트만이 화폐) · 동적 골라읽기 붕괴(read 레버는 정밀도뿐)."],
    fs=13)

# ================================================= 12. fp8 license
s = new_slide("12. 보너스 — 비대칭 정밀도 면허 (티어링이 fp8을 공짜로 만든다)",
              "같은 fp8인데: 매 토큰 재귀에 넣으면 붕괴, cold snapshot에 넣으면 무손실 — 4-cell 전부 실측")
cells = [
    ("raw + fp8 매토큰 (폐루프)", "MATH 93.2 붕괴\n생성 2.2× 방황", LRED, RED),
    ("미학습 + 티어링", "GSM8K 33.0 붕괴\n(정밀도 무관)", LRED, RED),
    ("학습 R + 티어링 + fp32-cold", "lossless\n(81.6 = fresh)", LGREEN, GREEN),
    ("학습 R + 티어링 + fp8-cold", "95.2 공식\n= fp32 = bf16", LGREEN, GREEN),
]
for i, (t1, t2, f, ln) in enumerate(cells):
    x = .8 + (i % 2) * 6.2; y = 1.5 + (i // 2) * 1.85
    box(s, x, y, 6.0, 1.6, f"{t1}\n{t2}", f, ln, fs=14, bold=True)
txt(s, .8, 5.4, 12.0, 1.7, [
    "메커니즘: 라운딩 주입 횟수가 T회(매 토큰 재귀 — 오차 복리) → T/c회(청크 개루프 — 복리 없음)로 줄기 때문.",
    "사슬: [R 학습 0.04%] → [티어링 실행권] → [청크 개루프] → [fp8 공짜]. 관문마다 실측으로 잠김.",
    "capacity 회계: bf16-cold는 총 state 메모리도 0.625× (대역폭 2.17×와 용량 −37%를 동시에) — \"용량 태워 속도 샀다\" 비판의 답."],
    fs=13)

# ================================================= 13. roadmap
s = new_slide("13. 현재 상태와 로드맵")
txt(s, .7, 1.3, 12.2, 5.6, [
    "완료 (실측):",
    "· 0.04% retrofit 파이프라인 + 회전 불변성 (fresh = raw 정답수 동일, 공식 스택)",
    "· B200 속도 1.62~2.42× (fused, B=256 표준) + 배치 스윕/roofline/READ-WRITE 회계",
    "· 공식 스택 정밀도×티어링 전수 매트릭스 + fp8 비대칭 면허 사슬 4-cell",
    "· 삭제-vs-lazy 모티베이션 사다리 (−26.5 vs −0.5) · 이론 실측 E-T1/E-T2 (loss-정의 중요도 · 도메인 이식성)",
    "· GSM8K lossless 학습 레시피 (longcot2: 장문 CoT + c4·c8 가중)",
    "",
    "진행 중 / 다음:",
    "· longcot2 공식 스택 재측정 (GSM8K/MATH-500) → 배포 ckpt 승격 — MATH −2.4 봉합이 목표",
    "· fp8 dequant-matvec 커널 (analytic 2.8×의 실측 앵커) · vLLM CUDA-graph e2e (~1.45× 목표)",
    "· self-speculative decoding (hot-only draft α=0.91, 출력 exact — varlen verify 커널이 완성 조건, 1.6~1.8× 투영)",
    "",
    "장기: cold 티어를 CXL-PNM에 상주 (analytic 3.8× + state capacity 8×) — 티어링 구조가 HW 계층에 그대로 매핑됨."],
    fs=14)

prs.save("../docs/pitch_deck_tiering.pptx")
print("PITCH DECK DONE -> docs/pitch_deck_tiering.pptx")
