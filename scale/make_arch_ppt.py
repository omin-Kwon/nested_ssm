"""Self-contained explainer deck: how Nemotron-9B computes, from attention
basics to our hot/cold tiered state. Figures drawn with matplotlib (English
labels), narration in Korean. Output: docs/nemotron_arch_explainer.pptx

Run: ~/nemo_env/bin/python3 make_arch_ppt.py   (from scale/)
"""
import json, os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mp
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TMP = "results/plots/ppt_figs"
os.makedirs(TMP, exist_ok=True)
FW, FH = 13.333, 7.5                      # 16:9 slide inches

def fig_path(name):
    return f"{TMP}/{name}.png"

def save(fig, name):
    fig.savefig(fig_path(name), dpi=170, bbox_inches="tight")
    plt.close(fig)

def box(ax, x, y, w, h, text, fc="#e8f0fe", ec="#1a56db", fs=9, tc="black"):
    ax.add_patch(mp.FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02",
                                   fc=fc, ec=ec, lw=1.2))
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center",
            fontsize=fs, color=tc)

def arrow(ax, x0, y0, x1, y1, **kw):
    ax.annotate("", xy=(x1, y1), xytext=(x0, y0),
                arrowprops=dict(arrowstyle="->", lw=1.4,
                                color=kw.get("color", "#333")))

# ---------------------------------------------------------------- figures
# A. KV growth vs fixed state
fig, axes = plt.subplots(1, 2, figsize=(9, 3.4))
ax = axes[0]
for t in range(1, 9):
    ax.add_patch(mp.Rectangle((t - 1, 0), .9, t, fc="#fca5a5", ec="k", lw=.4))
ax.set_xlim(-.2, 8.4); ax.set_ylim(0, 9.5)
ax.set_title("Transformer: KV cache grows with T", fontsize=10)
ax.set_xlabel("decode step t"); ax.set_ylabel("memory")
ax.text(4, 8.7, "attend to ALL past tokens:  O(T) memory, O(T) work/token",
        ha="center", fontsize=8)
ax = axes[1]
for t in range(1, 9):
    ax.add_patch(mp.Rectangle((t - 1, 0), .9, 3, fc="#93c5fd", ec="k", lw=.4))
ax.set_xlim(-.2, 8.4); ax.set_ylim(0, 9.5)
ax.set_title("Linear attention / SSM: state is FIXED size", fontsize=10)
ax.set_xlabel("decode step t")
ax.text(4, 8.7, "past is COMPRESSED into a fixed matrix S:  O(1) memory & work",
        ha="center", fontsize=8)
for a in axes: a.set_xticks([]); a.set_yticks([])
save(fig, "A_kv_vs_state")

# B. from softmax attention to accumulator
fig, ax = plt.subplots(figsize=(9.5, 3.2)); ax.axis("off")
ax.text(.5, .92, "softmax attention:", fontsize=11, transform=ax.transAxes,
        ha="right", color="#b91c1c")
ax.text(.52, .92, r"$y_t=\sum_{s\leq t}\ \mathrm{softmax}(q_t\!\cdot\!k_s)\,v_s$"
        "   — every step touches every past token",
        fontsize=12, transform=ax.transAxes)
ax.text(.5, .60, "drop softmax:", fontsize=11, transform=ax.transAxes,
        ha="right", color="#1d4ed8")
ax.text(.52, .60, r"$y_t=\sum_{s\leq t}(q_t\!\cdot\!k_s)\,v_s"
        r"=q_t\cdot S_t$,   $S_t \equiv \sum_{s\leq t}k_s v_s^\top$",
        fontsize=13, transform=ax.transAxes)
ax.text(.52, .24, r"$S_t = S_{t-1} + k_t v_t^\top$   →  past summarized by ONE"
        " matrix, updated incrementally", fontsize=12, transform=ax.transAxes)
save(fig, "B_softmax_to_state")

# C. the state object
fig, ax = plt.subplots(figsize=(9, 3.6)); ax.axis("off")
ax.add_patch(mp.Rectangle((.05, .15), .35, .6, fc="#dbeafe", ec="#1d4ed8", lw=1.5))
ax.text(.225, .45, "S\n(P × N)", ha="center", va="center", fontsize=14)
ax.text(.225, .08, "one head's memory", ha="center", fontsize=9)
ax.annotate("P = head_dim = 64\n(value channels)", xy=(.05, .45),
            xytext=(-.02, .45), ha="right", fontsize=9, va="center")
ax.annotate("N = d_state = 128  (key/slot axis — B,C index this)",
            xy=(.225, .78), ha="center", fontsize=9)
ax.text(.55, .62, "Nemotron-Nano-9B scale:", fontsize=11, fontweight="bold")
ax.text(.55, .50, "128 heads × (64×128) per mamba layer", fontsize=10)
ax.text(.55, .40, "27 mamba layers  →  113 MB fp32 per REQUEST", fontsize=10)
ax.text(.55, .30, "B=256 requests  →  29 GB live state", fontsize=10)
ax.text(.55, .16, "every decode step must READ and WRITE all of it",
        fontsize=10, color="#b91c1c")
save(fig, "C_state_object")

# D. one decode step pipeline
fig, ax = plt.subplots(figsize=(11, 3.6)); ax.axis("off")
ax.set_xlim(0, 11); ax.set_ylim(0, 3.6)
box(ax, .1, 1.4, 1.15, .8, "x_t\n(b, d=4480)", fc="#f3f4f6", ec="#555")
box(ax, 1.55, 1.4, 1.3, .8, "in_proj\n(GEMM)")
box(ax, 3.15, 1.4, 1.5, .8, "causal conv1d\n+ SiLU  (window 4)")
box(ax, 4.95, 1.4, 1.7, .8, "split →\nh (b,8192)\nB,C (b,8×128)\ndt (b,128)", fs=8)
box(ax, 6.95, 1.4, 1.9, .8, "SSM state update\n+ readout  (SSU)",
    fc="#fee2e2", ec="#b91c1c")
box(ax, 9.15, 1.4, 1.0, .8, "norm ×\ngate")
box(ax, 10.25, 1.4, .68, .8, "out\nproj")
for x0, x1 in [(1.25, 1.55), (2.85, 3.15), (4.65, 4.95), (6.65, 6.95),
               (8.85, 9.15), (10.15, 10.25)]:
    arrow(ax, x0, 1.8, x1, 1.8)
box(ax, 6.95, .2, 1.9, .75, "state S\n27L × (128H,64P,128N)",
    fc="#dbeafe", ec="#1d4ed8", fs=8)
arrow(ax, 7.6, .95, 7.6, 1.4, color="#1d4ed8")
arrow(ax, 8.2, 1.4, 8.2, .95, color="#b91c1c")
ax.text(7.45, 1.13, "read", fontsize=8, color="#1d4ed8", ha="right")
ax.text(8.32, 1.13, "write", fontsize=8, color="#b91c1c")
ax.text(5.5, 3.25, "one mamba2 layer, one decode step (per token)",
        fontsize=11, ha="center", fontweight="bold")
ax.text(6.95, 2.45, "gate (h-path)", fontsize=8, color="#555")
save(fig, "D_pipeline")

# E. update geometry: decay + outer product
fig, ax = plt.subplots(figsize=(10, 3.4)); ax.axis("off")
ax.set_xlim(0, 10); ax.set_ylim(0, 3.4)
ax.add_patch(mp.Rectangle((.3, .6), 1.8, 1.8, fc="#dbeafe", ec="#1d4ed8"))
ax.text(1.2, 1.5, r"$S_{t-1}$", fontsize=13, ha="center")
ax.text(2.55, 1.5, r"$\times\,a_t$", fontsize=12, ha="center")
ax.text(2.55, 1.1, "decay\n(scalar/head)", fontsize=8, ha="center")
ax.text(3.15, 1.5, "+", fontsize=15, ha="center")
ax.add_patch(mp.Rectangle((3.6, .6), .22, 1.8, fc="#fecaca", ec="#b91c1c"))
ax.text(3.71, 2.6, r"$\Delta x_t$ (P)", fontsize=9, ha="center")
ax.text(4.15, 1.5, r"$\otimes$", fontsize=14, ha="center")
ax.add_patch(mp.Rectangle((4.5, 1.4), 1.8, .22, fc="#fed7aa", ec="#c2410c"))
ax.text(5.4, 1.9, r"$B_t$ (N)", fontsize=9, ha="center")
ax.text(6.6, 1.5, "=", fontsize=15, ha="center")
ax.add_patch(mp.Rectangle((7.0, .6), 1.8, 1.8, fc="#dbeafe", ec="#1d4ed8"))
ax.add_patch(mp.Rectangle((7.0, .6), 1.8, 1.8, fill=False, hatch="..",
                           ec="#b91c1c", lw=0))
ax.text(7.9, 1.5, r"$S_t$", fontsize=13, ha="center")
ax.text(7.9, .25, "rank-1 'stamp' of the new token pressed onto the whole sheet",
        fontsize=8, ha="center")
ax.text(5, 3.1, r"update:  $S_t = a_t\,S_{t-1} + (\Delta x_t)\otimes B_t$"
        r"     readout:  $y_t = S_t\,C_t + D\odot x_t$",
        fontsize=13, ha="center")
save(fig, "E_update_geometry")

# F. bytes per step (measured)
R = {int(k): v for k, v in
     json.load(open("results/vllm_sweep_breakdown.json")).items()}
BS = [B for B in sorted(R) if R[B].get("state_dtype", "float32") == "float32"]
fig, ax = plt.subplots(figsize=(8.5, 3.6))
ssu = [R[B]["buckets_ms"].get("state-op (SSU)", 0) / R[B]["gen"] for B in BS]
gem = [R[B]["buckets_ms"].get("GEMM (proj/MLP/lm_head)", 0) / R[B]["gen"] for B in BS]
rest = [R[B]["ms_per_step_busy"] - s - g for B, s, g in zip(BS, ssu, gem)]
x = np.arange(len(BS))
ax.bar(x, [s / 2 for s in ssu], .6, color="#d62728", label="state READ (+readout)")
ax.bar(x, [s / 2 for s in ssu], .6, bottom=[s / 2 for s in ssu],
       color="#d62728", hatch="\\\\", edgecolor="white", label="state WRITE")
ax.bar(x, gem, .6, bottom=ssu, color="#1f77b4", label="GEMM (weights)")
ax.bar(x, rest, .6, bottom=[s + g for s, g in zip(ssu, gem)], color="#bbb",
       label="everything else")
ax.set_xticks(x, [str(b) for b in BS])
ax.set_xlabel("batch size"); ax.set_ylabel("ms / decode step (GPU busy)")
ax.set_title("Measured on B200 (vLLM): state movement becomes THE cost", fontsize=10)
ax.legend(fontsize=8); ax.grid(axis="y", alpha=.25)
save(fig, "F_measured")

# G. hybrid stack
fig, ax = plt.subplots(figsize=(10.5, 1.9)); ax.axis("off")
pat = "M-M-M-MM-M-M-M*-M-M-M*-M-M-M-M*-M-M-M-M*-M-MM-M-M-M-M-M-"
cols = {"M": "#1d4ed8", "*": "#b91c1c", "-": "#d1d5db"}
for i, ch in enumerate(pat):
    ax.add_patch(mp.Rectangle((i * .19, .35), .17, .9, fc=cols[ch], ec="none"))
ax.set_xlim(-.1, len(pat) * .19 + .1); ax.set_ylim(0, 1.9)
ax.text(0, 1.55, "Nemotron-Nano-9B-v2: 56 layers = 27 Mamba (blue) + 4 attention"
        " (red) + 25 MLP (gray)", fontsize=10)
ax.text(0, .06, "sequence memory lives almost entirely in the 27 mamba states"
        " -> that's what we tier", fontsize=9, color="#1d4ed8")
save(fig, "G_hybrid")

# H. hot/cold tiering
fig, ax = plt.subplots(figsize=(10.5, 4.2)); ax.axis("off")
ax.set_xlim(0, 10.5); ax.set_ylim(0, 4.2)
ax.add_patch(mp.Rectangle((.4, 1.2), .9, 2.2, fc="#fecaca", ec="#b91c1c", lw=1.5))
ax.add_patch(mp.Rectangle((1.3, 1.2), 2.7, 2.2, fc="#bfdbfe", ec="#1d4ed8", lw=1.5))
ax.text(.85, 3.55, "hot\npb=32", ha="center", fontsize=10, color="#b91c1c")
ax.text(2.65, 3.55, "cold  (96 of N=128 columns)", ha="center", fontsize=10,
        color="#1d4ed8")
ax.text(.85, 2.3, "fresh\nevery\ntoken", ha="center", fontsize=9)
ax.text(2.65, 2.3, "stale snapshot (bf16/fp8)\nread every token\n"
        "decay-compensated ×exp(glog)", ha="center", fontsize=9)
ax.text(2.2, .9, "N axis (key/slot dimension) — split HERE", ha="center",
        fontsize=9)
box(ax, 5.0, 2.6, 2.3, .9,
    "every token:\n y = C·S[:, :pb]  (hot, fresh)\n   + C·snap × exp(Σ dt·A)",
    fc="#fff7ed", ec="#c2410c", fs=9)
box(ax, 5.0, 1.0, 2.3, .9,
    "every c tokens (FLUSH):\n fold c updates EXACTLY\n snap ← quantize(cold)",
    fc="#ecfdf5", ec="#047857", fs=9)
box(ax, 7.9, 1.8, 2.4, 1.2,
    "traffic per step (vs raw=2.0)\nhot R+W 0.5\ncold read 0.19 (fp8)\n"
    "flush ~0.94/c\n→ c16: 2.68× less state traffic", fc="#f3f4f6", ec="#555", fs=9)
ax.text(5.25, 3.9, "the deal: cold updates are DELAYED (exact), cold reads are"
        " STALE (approximate, decay-corrected)", fontsize=10)
save(fig, "H_tiering")

# I. rotation
fig, ax = plt.subplots(figsize=(10, 3.4)); ax.axis("off")
ax.set_xlim(0, 10); ax.set_ylim(0, 3.4)
np.random.seed(0)
val = np.abs(np.random.randn(6, 12)) * np.linspace(1, 1, 12)
ax.imshow(val, extent=(.4, 4.4, .8, 2.8), cmap="Reds", aspect="auto")
ax.text(2.4, 3.0, "pretrained: information scattered across N columns",
        fontsize=9, ha="center")
val2 = np.sort(val, axis=1)[:, ::-1]
ax.imshow(val2, extent=(5.6, 9.6, .8, 2.8), cmap="Reds", aspect="auto")
ax.text(7.6, 3.0, "after learning R: concentrated into the first (hot) columns",
        fontsize=9, ha="center")
arrow(ax, 4.55, 1.8, 5.45, 1.8)
ax.text(5.0, 2.05, "R", fontsize=13, ha="center", color="#1d4ed8")
ax.text(5.0, .35, r"$B\!\to\!RB,\ C\!\to\!RC$ (orthogonal, per group):"
        r"  $y=S R^\top R\, C = S\,C$  — output EXACTLY unchanged at full width",
        fontsize=10, ha="center")
save(fig, "I_rotation")

# J. throughput takeaway
fig, ax = plt.subplots(figsize=(8.5, 3.6))
tp = [R[B]["tok_per_s"] for B in BS]
ax.plot(BS, tp, "o-", lw=2, color="#1f77b4", label="raw (measured)")
for cut, lab, c in [(2.68, "v4-c16-fp8 (projected)", "#d62728")]:
    tp2 = []
    for B in BS:
        s = R[B]["buckets_ms"].get("state-op (SSU)", 0) / R[B]["gen"]
        tp2.append(B * 1e3 / (R[B]["ms_per_step_wall"] - s * (1 - 1 / cut)))
    ax.plot(BS, tp2, "s--", lw=1.4, color=c, label=lab)
ax.set_xscale("log", base=2); ax.set_xticks(BS, [str(b) for b in BS], fontsize=8)
ax.set_xlabel("batch size"); ax.set_ylabel("decode tok/s")
ax.set_title("The ceiling is set by state bytes — cut the bytes, raise the ceiling",
             fontsize=10)
ax.legend(fontsize=9); ax.grid(alpha=.25)
save(fig, "J_takeaway")

# ---------------------------------------------------------------- deck
prs = Presentation()
prs.slide_width = Inches(FW); prs.slide_height = Inches(FH)
BLANK = prs.slide_layouts[6]

def slide(title, bullets=None, img=None, img_w=None, top=1.25, tfs=15):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(.45), Inches(.2), Inches(FW - .9), Inches(.85))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(27); p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x36, 0x5d)
    y = top
    if bullets:
        bb = s.shapes.add_textbox(Inches(.55), Inches(y), Inches(FW - 1.1),
                                  Inches(1.8))
        tf = bb.text_frame; tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b; p.font.size = Pt(tfs)
        y += .42 * len(bullets) + .25
    if img:
        from PIL import Image
        iw, ih = Image.open(fig_path(img)).size
        w = img_w or min(FW - 1.2, (FH - y - .25) * iw / ih)
        h = w * ih / iw
        if y + h > FH - .15:
            h = FH - .15 - y; w = h * iw / ih
        s.shapes.add_picture(fig_path(img), Inches((FW - w) / 2), Inches(y),
                             Inches(w), Inches(h))
    return s

# 1 title
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11.3), Inches(2.5))
tf = tb.text_frame
tf.paragraphs[0].text = "Nemotron-9B은 어떻게 계산하는가"
tf.paragraphs[0].font.size = Pt(40); tf.paragraphs[0].font.bold = True
p = tf.add_paragraph(); p.text = ("attention의 문제에서 recurrent state까지, "
                                  "그리고 우리의 hot/cold tiered state까지 — self-contained 해설")
p.font.size = Pt(18)
p = tf.add_paragraph(); p.text = "Elastic Test-Time Memory 프로젝트 (2026-07)"
p.font.size = Pt(14); p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# 2
slide("1. 왜 Transformer 그대로는 안 되는가",
      ["Transformer decode는 매 토큰마다 과거 '전부'(KV cache)를 읽어야 함 — 길이 T에 비례해 메모리·연산 증가",
       "Linear attention/SSM 계열: 과거를 고정 크기 행렬 S로 '압축'해 들고 다님 — 토큰당 비용이 T와 무관"],
      "A_kv_vs_state")
# 3
slide("2. softmax를 지우면 '누적 상태'가 나온다",
      ["softmax만 제거하면 합의 순서를 바꿀 수 있고, 과거 전체가 하나의 행렬 S로 접힘",
       "S는 '지금까지 본 모든 (key, value) 쌍의 요약'이고, 매 토큰 덧셈 한 번으로 갱신됨"],
      "B_softmax_to_state")
# 4
slide("3. 상태 S의 실체 — Nemotron 스케일",
      ["head 하나의 기억 = P×N 행렬 (P=값 채널 64, N=key/슬롯 축 128) — B와 C가 N축을 인덱싱",
       "요청 1개당 113MB(fp32), B=256이면 29GB가 '살아있는 메모리' — 매 스텝 전량 R/W"],
      "C_state_object")
# 5
slide("4. mamba2 레이어의 decode 한 스텝 (파이프라인)",
      ["GEMM(projection) → 4-tap causal conv+SiLU → (h, B, C, dt)로 분해 → SSM 갱신+readout → gate·norm → out_proj",
       "빨간 박스(SSU)가 상태를 읽고 쓰는 유일한 지점 — 여기가 대역폭의 전쟁터"],
      "D_pipeline")
# 6
slide("5. 상태 갱신의 기하학 — '도장 찍기 + 바래기'",
      ["갱신 = (i) 기존 기억 전체를 스칼라 a_t로 살짝 바래게 하고(decay), (ii) 새 토큰의 rank-1 도장(Δx⊗B)을 찍는 것",
       "readout = 질의 C로 N축을 따라 내적 — '어느 슬롯의 기억을 얼마나 꺼낼까'"],
      "E_update_geometry")
# 7
slide("6. 실측: 배치가 커지면 상태 이동이 비용의 본체",
      ["B200 실측(vLLM): B=256에서 state R/W가 GPU 시간의 58% — GEMM(가중치)은 배치가 공유해 거의 고정",
       "SSU 커널은 4.6TB/s로 대역폭 포화 — 시간 ∝ 바이트, read와 write가 정확히 절반씩"],
      "F_measured")
# 8
slide("7. Nemotron-H 하이브리드 구조",
      ["56층 중 attention은 4층뿐 — 시퀀스 기억의 본체는 27개 mamba 상태",
       "그래서 mamba 상태만 다뤄도 decode 비용의 지배항을 다루는 것"],
      "G_hybrid")
# 9
slide("8. 우리의 v4: N축을 hot/cold로 가른다",
      ["hot(앞 32열)은 매토큰 fresh, cold(96열)는 c토큰마다 정확히 몰아쓰기(flush) + 그 사이엔 stale 스냅샷 읽기",
       "쓰기는 '빈도'로, 읽기는 '바이트(bf16/fp8)'로 줄인다 — c16-fp8이면 상태 트래픽 2.68×↓"],
      "H_tiering")
# 10
slide("9. 회전 R — 모델을 안 바꾸고 서랍만 재배치",
      ["B·C를 같은 직교 R로 돌리면 출력이 '정확히' 불변 (RᵀR=I) — full width에선 수학적 무손실",
       "학습되는 것은 R(0.04%)뿐: 중요한 정보를 hot 열로 모으는 좌표 재배치 — 실측 fresh=raw 동률(GSM8K 95.0)"],
      "I_rotation")
# 11
slide("10. 결과: 천장을 들어올리다",
      ["decode throughput의 천장은 state 바이트가 결정 (B≥512에서 11.9k tok/s 포화 실측)",
       "v4는 그 바이트 자체를 줄여 천장을 올림 — 정확도는 3-arm(raw/fresh/v4)으로 lossless 검증"],
      "J_takeaway")
# 12 summary equations
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(.45), Inches(.2), Inches(12.4), Inches(.85))
p = tb.text_frame.paragraphs[0]; p.text = "한 장 요약 — 수식 카드"
p.font.size = Pt(27); p.font.bold = True
p.font.color.rgb = RGBColor(0x1a, 0x36, 0x5d)
body = s.shapes.add_textbox(Inches(.7), Inches(1.15), Inches(12), Inches(5.9))
tf = body.text_frame; tf.word_wrap = True
rows = [
 ("기본 갱신", "S_t = a_t · S_{t-1} + (Δx_t) ⊗ B_t     y_t = S_t · C_t + D⊙x_t"),
 ("회전 retrofit", "B→RB, C→RC (직교, 그룹별)  ⇒  S→SRᵀ, y 불변 — hot 정렬은 학습으로"),
 ("v4 readout", "y = C_hot·S[:, :pb]  +  C_cold·snap × exp(Σ_{since flush} dt·A)"),
 ("v4 flush (정확)", "c토큰마다: S_cold ← 접힌(fold) 누적 갱신, snap ← quant(S_cold)"),
 ("정밀도 면허", "라운딩 주입 T회(raw-fp8, 폐루프: 붕괴) vs T/c회(cold-fp8, 개루프: 무손실)"),
 ("트래픽 (raw=2.0)", "hot 0.5 + cold_read 0.19(fp8) + flush 0.94/c  →  c16: 2.68×↓"),
]
for i, (k, v) in enumerate(rows):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{k}:   {v}"
    p.font.size = Pt(16)
    p.space_after = Pt(14)

prs.save("../docs/nemotron_arch_explainer.pptx")
print("PPT DONE -> docs/nemotron_arch_explainer.pptx")
