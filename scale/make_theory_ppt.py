"""Theory deck: why one learned rotation sorts recurrent memory & generalizes.
Visualizes the 2026-07-14 theory discussion. Output: docs/theory_deck.pptx
Run: ~/nemo_env/bin/python3 make_theory_ppt.py   (from scale/)
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mp
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TMP = "results/plots/theory_figs"
os.makedirs(TMP, exist_ok=True)
FW, FH = 13.333, 7.5

def save(fig, name):
    fig.savefig(f"{TMP}/{name}.png", dpi=170, bbox_inches="tight")
    plt.close(fig)

# ---- F1: the cloud of B/C vectors with principal directions ----
rng = np.random.default_rng(0)
n = 600
t = rng.choice(3, n, p=[.6, .25, .15])
dirs = np.array([[1, 0.35], [-0.4, 1], [0.9, -0.9]])
dirs = dirs / np.linalg.norm(dirs, axis=1, keepdims=True)
pts = dirs[t] * rng.normal(1.2, .35, (n, 1)) + rng.normal(0, .12, (n, 2))
fig, ax = plt.subplots(figsize=(6.4, 5.4))
cols = np.array(["#1d4ed8", "#c2410c", "#047857"])
ax.scatter(pts[:, 0], pts[:, 1], s=8, c=cols[t], alpha=.45)
for d, c, lab, p in zip(dirs, cols, ["u1 (60%)", "u2 (25%)", "u3 (15%)"],
                        [.6, .25, .15]):
    ax.annotate("", xy=d * 2.3, xytext=-d * 2.3,
                arrowprops=dict(arrowstyle="<->", lw=1 + 6 * p, color=c))
    ax.annotate(lab, xy=d * 2.35, fontsize=11, color=c, fontweight="bold")
ax.set_xlim(-2.9, 2.9); ax.set_ylim(-2.9, 2.9)
ax.set_xticks([]); ax.set_yticks([])
ax.set_title("the CLOUD of write/read vectors (B, C) over the corpus\n"
             "each token = one point; traffic directions = cloud's principal axes",
             fontsize=10)
ax.axhline(0, color="#999", lw=.5); ax.axvline(0, color="#999", lw=.5)
ax.annotate("raw axes (drawers 1,2):\nmisaligned with the cloud",
            xy=(2.0, 0.12), fontsize=8, color="#666")
save(fig, "F1_cloud")

# ---- F2: traffic bars — raw flat vs rotated sorted ----
fig, axes = plt.subplots(1, 2, figsize=(9.2, 3.4))
raw = [0.22, 0.22, 0.12, 0.22, 0.22]
rot = [0.60, 0.25, 0.15, 0.0, 0.0]
for ax, v, ttl in [(axes[0], raw, "raw axes: T = (0.22, 0.22, 0.12, 0.22, 0.22)\n"
                    "concentration EXISTS but is invisible"),
                   (axes[1], rot, "after R: T' = eigenvalues of M = (0.60, 0.25, 0.15, 0, 0)\n"
                    "hot = drawers 1-2 covers 85%")]:
    colors = ["#d62728" if i < 2 else "#1f77b4" for i in range(5)]
    ax.bar(range(1, 6), v, color=colors)
    ax.set_ylim(0, .7); ax.set_xlabel("drawer n"); ax.set_title(ttl, fontsize=9)
    ax.set_ylabel("avg traffic share")
axes[1].axvline(2.5, color="k", ls="--", lw=1)
axes[1].annotate("hot | cold", xy=(2.55, .62), fontsize=10)
save(fig, "F2_traffic")

# ---- F3: passive rotation — same arrow, new axes ----
fig, ax = plt.subplots(figsize=(5.6, 5))
ax.annotate("", xy=(1, 1), xytext=(0, 0),
            arrowprops=dict(arrowstyle="->", lw=3, color="#111"))
ax.annotate("x = (1,1)  [raw axes]", xy=(1.02, 1.02), fontsize=10)
ax.axhline(0, color="#1f77b4", lw=1.5); ax.axvline(0, color="#1f77b4", lw=1.5)
for d, lab in [((1, 1), "r1 = u"), ((1, -1), "r2 = v")]:
    d = np.array(d) / np.sqrt(2)
    ax.annotate("", xy=d * 1.9, xytext=-d * 1.9,
                arrowprops=dict(arrowstyle="->", lw=1.2, ls="--", color="#d62728"))
    ax.annotate(lab, xy=d * 1.95, color="#d62728", fontsize=10)
ax.text(-1.85, 1.7, "R·x = (√2, 0):\nSAME arrow,\nre-read on new axes\n"
        "(RᵀR = I keeps all\nlengths & angles)", fontsize=10,
        bbox=dict(fc="#fff7ed", ec="#c2410c"))
ax.set_xlim(-2, 2); ax.set_ylim(-2, 2); ax.set_xticks([]); ax.set_yticks([])
ax.set_title("multiplying by an ORTHOGONAL R = passive coordinate change",
             fontsize=10)
save(fig, "F3_passive")

# ---- F4: SGD as voting / streaming eigendecomposition ----
fig, ax = plt.subplots(figsize=(8.6, 3.6))
steps = np.arange(0, 400)
rng2 = np.random.default_rng(1)
align1 = 1 - 0.9 * np.exp(-steps / 80) + rng2.normal(0, .01, 400)
align2 = 1 - 0.95 * np.exp(-steps / 140) + rng2.normal(0, .012, 400)
align3 = 1 - 0.97 * np.exp(-steps / 230) + rng2.normal(0, .015, 400)
for a, c, lab in [(align1, "#1d4ed8", "row 1 -> u1 (60% of votes)"),
                  (align2, "#c2410c", "row 2 -> u2 (25%)"),
                  (align3, "#047857", "row 3 -> u3 (15%)")]:
    ax.plot(steps, np.clip(a, 0, 1), color=c, label=lab)
ax.set_xlabel("training step (each sample = one loss-weighted 'vote')")
ax.set_ylabel("alignment |r_n · u_n|")
ax.set_title("SGD never forms M — it tallies votes and converges to M's eigenbasis\n"
             "(nested dropout = PCA theorem; streaming-PCA / Oja's-rule family)",
             fontsize=9)
ax.legend(fontsize=8); ax.grid(alpha=.3)
save(fig, "F4_voting")

# ---- F5: generalization — MEASURED (E-T1 refutation + E-T2 transfer) ----
import json as _json
_et1 = _json.load(open("results/M_spectrum.json"))["mean"]
_et1b = _json.load(open("results/M_lagged.json"))
_et2 = _json.load(open("results/domain_elasticity.json"))
fig, axes = plt.subplots(1, 2, figsize=(10.2, 3.7))
# (a) E-T1/E-T1b: trained R vs explicit-statistic eigenbases — 4-way convergence
names = ["explicit M\n(lag-0)", "lagged M\n(a=0.9)", "lagged M\n(a=0.99)"]
cosv = [_et1["subspace_cos"], _et1b["0.9"]["cos"], _et1b["0.99"]["cos"]]
axes[0].bar(names, cosv, color="#1f77b4", width=.55)
axes[0].axhline(0.5, color="#d62728", ls="--", lw=1.5)
axes[0].annotate("random 32-dim subspace (~0.5)", xy=(0.02, 0.51),
                 color="#d62728", fontsize=8)
axes[0].axhline(1.0, color="#047857", ls=":", lw=1.2)
axes[0].annotate("perfect match (1.0)", xy=(1.6, 0.93), color="#047857", fontsize=8)
axes[0].set_ylim(0, 1.05); axes[0].set_ylabel("principal-angle cos (trained R vs eigenbasis)")
axes[0].set_title("MEASURED: trained hot subspace is UNCORRELATED with any\n"
                  f"training-free statistic's eigenbasis (gap@32 = {_et1['gap_at_32']:.4f} too)",
                  fontsize=9)
# (b) E-T2: truncation-tax curves per domain × config — only trained R transfers
ws = [8, 16, 32, 64, 96, 128]
dcol = {"wiki": "#1d4ed8", "math": "#c2410c", "code": "#047857"}
styles = {"trained": ("o-", 2.2), "ghost": ("s--", 1.2), "identity": ("^:", 1.0)}
for cfg, (st, lw) in styles.items():
    for dom, c in dcol.items():
        tax = [_et2[cfg][dom]["tax"][str(k)] for k in ws]
        lab = f"{cfg} R" if dom == "wiki" else None
        axes[1].plot(ws, tax, st, color=c, lw=lw, ms=4, alpha=1 if cfg == "trained" else .6,
                     label=lab)
axes[1].axhline(1.0, color="#999", lw=.8)
axes[1].set_xscale("log", base=2); axes[1].set_yscale("log")
axes[1].set_xticks(ws); axes[1].set_xticklabels(ws)
axes[1].set_xlabel("hot width k"); axes[1].set_ylabel("ppl tax  ppl(k)/ppl(128)  [log]")
axes[1].set_title("MEASURED: ONE trained R -> near-1.0 tax on wiki/math/code alike;\n"
                  "identity & calibration(GHOST) bases pay 1.6-400x, domain-erratic",
                  fontsize=9)
axes[1].legend(fontsize=8); axes[1].grid(alpha=.3, which="both")
save(fig, "F5_measured")

# ---- F6: union-subspace condition ----
fig, ax = plt.subplots(figsize=(7.4, 3.8)); ax.axis("off")
ax.set_xlim(0, 10); ax.set_ylim(0, 4)
ax.add_patch(mp.FancyBboxPatch((0.4, 0.6), 5.6, 2.8, boxstyle="round,pad=0.05",
                               fc="#fee2e2", ec="#b91c1c", lw=2))
ax.text(3.2, 3.6, "hot = a 32-dim SUBSPACE (order inside is irrelevant)",
        ha="center", fontsize=10, color="#b91c1c")
for cx, cy, r, c, lab in [(2.0, 2.0, 1.15, "#1d4ed8", "task A\ntop dirs"),
                          (3.4, 2.0, 1.15, "#047857", "task B\ntop dirs"),
                          (4.6, 2.0, 1.15, "#c2410c", "task C\ntop dirs")]:
    ax.add_patch(mp.Circle((cx, cy), r, fill=False, ec=c, lw=2))
    ax.text(cx, cy + r + .12, lab, ha="center", fontsize=8, color=c)
ax.text(7.9, 2.6, "requirement is NOT\n'same eigen-ORDER':\n\n"
        "union of tasks' top\ndirections fits in the\n32-dim hot subspace",
        fontsize=10, va="center")
save(fig, "F6_union")

# ---------------------------------------------------------------- deck
prs = Presentation()
prs.slide_width = Inches(FW); prs.slide_height = Inches(FH)
BLANK = prs.slide_layouts[6]

def slide(title, bullets=None, img=None, top=1.2, tfs=14):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(.45), Inches(.2), Inches(FW - .9), Inches(.85))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(26); p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x36, 0x5d)
    y = top
    if bullets:
        bb = s.shapes.add_textbox(Inches(.55), Inches(y), Inches(FW - 1.1), Inches(2))
        tf = bb.text_frame; tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b; p.font.size = Pt(tfs)
        y += .40 * len(bullets) + .25
    if img:
        from PIL import Image
        iw, ih = Image.open(f"{TMP}/{img}.png").size
        w = min(FW - 1.4, (FH - y - .25) * iw / ih)
        h = w * ih / iw
        if y + h > FH - .15:
            h = FH - .15 - y; w = h * iw / ih
        s.shapes.add_picture(f"{TMP}/{img}.png", Inches((FW - w) / 2), Inches(y),
                             Inches(w), Inches(h))
    return s

# 1 title
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(.9), Inches(2.3), Inches(11.6), Inches(2.6))
tf = tb.text_frame
tf.paragraphs[0].text = "왜 '하나의 직교 회전'이 기억을 정렬하고, 왜 일반화되는가"
tf.paragraphs[0].font.size = Pt(34); tf.paragraphs[0].font.bold = True
p = tf.add_paragraph(); p.text = ("N차원 key 공간의 기저 재선택으로서의 R — 통신 통계, "
                                  "loss-정의 중요도(E-T1), 도메인 이식 실측(E-T2)")
p.font.size = Pt(17)
p = tf.add_paragraph(); p.text = "Elastic Test-Time Memory · Theory Deck (2026-07-14)"
p.font.size = Pt(13); p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# 2 cloud
slide("1. 재료: 말뭉치 전체가 만드는 '벡터 구름'",
      ["토큰(문맥)마다 동결된 가중치가 key B, query C를 생성 — 말뭉치 전체로 보면 N차원 공간의 점 구름",
       "구름의 주축들 = '통신 방향' — 개별 토큰이 아니라 분포의 성질 (아래: 60/25/15 장난감)"],
      "F1_cloud")
# 3 traffic
slide("2. 집중은 이미 있었다 — 보이지 않았을 뿐",
      ["평균 통신량 행렬 M = E[C Bᵀ·w]: 방향 r의 트래픽 = rᵀMr — 앞 k서랍 최대화의 해 = M의 고유기저(고유값 순)",
       "원래 좌표에선 트래픽이 평평(왼쪽) → 회전 후 고유값 수열로 정렬(오른쪽) — R은 집중을 '창조'하지 않고 '축에 정렬'"],
      "F2_traffic")
# 4 orthogonal
slide("3. 직교행렬 = 무손실의 계약서",
      ["정의 RᵀR = I (행들이 상호 수직·길이 1) → 내적·길이·각도 보존",
       "수동적 관점: 벡터는 그대로, 좌표계만 교체 — 모델의 계산(⟨C,B⟩)은 좌표 무관 → full-width 정확 불변 (fresh=raw 실증)"],
      "F3_passive")
# 5 SGD
slide("4. 학습 = M을 만들지 않는 고유분해",
      ["M은 (동결 가중치 × 데이터 분포)가 정하는 수학적 사실 — 아무도 계산하지 않음",
       "prefix-loss SGD: 표본마다 자기 방향에 loss-가중 '한 표' → 개표 결과가 고유값 순 (nested dropout=PCA 정리; Oja/streaming-PCA 계보)"],
      "F4_voting")
# 6 generalization core — REWIRED after E-T1/E-T2 measurements (07-14)
slide("5. 왜 과적합되지 않는가 (1) — 갭 논증이 아니라 '작은 가설류 + 이식 실측'",
      ["R의 가설 공간 = 직교군 (3.5M, 백본의 0.04%) — 특징을 새로 배우지 않고 기존 통계의 좌표만 고름",
       "이상화(Davis-Kahan): 갭 크면 부분공간 안정 — 실측이 기각: 명시 M은 gap@32≈0, 학습 R과 주각 cos 0.44≈무작위"
       " (lag-0·시차·열노름·질의에너지 4중 수렴) → 중요도는 loss가 정의, training-free 통계로 대체 불가",
       "일반화의 실제 근거 = ① 가설류 작음 ② 이식의 직접 실측(E-T2): 한 R의 절단세 곡선이 wiki/수학/코드에서 겹침"
       " (k16: 1.26/1.31/1.27×) vs calibration 기저는 k16에서 17-36× 절벽"],
      "F5_measured", tfs=13)
# 7 union
slide("6. 왜 과적합되지 않는가 (2) — 요구 조건이 '순서'가 아니라 '합집합'이라서",
      ["hot은 방향 1개가 아닌 그룹당 32차원 부분공간 — hot 내부 순서는 무의미, 경계만 유효",
       "태스크들이 서로 다른 순서로 밀어도 상위 방향 합집합이 32차원에 들어오면 무충돌 + 27층×8그룹이 다양성 분산 흡수",
       "실증: 일반 LM loss로만 학습한 R이 GSM8K/RULER/recall/HumanEval 전반 lossless — 상위 스펙트럼 = 태스크 공유 인프라",
       "정밀화(E-T2): 배우는 것은 전순서도 아닌 '메뉴 격자의 부분공간 사슬(flag)' — 메뉴 밖 폭(k=8, 96)은 3도메인"
       " 공통으로 유료(k96 1.2-1.4× > k64) → 학습이 굽는 건 메뉴 지점의 경계들뿐"],
      "F6_union", tfs=13)
# 8 boundary
slide("7. 경계 부근만 혼합에 민감 — 레시피의 통역",
      ["상위부(범용 채널)는 태스크 불변, hot/cold 경계 부근은 혼합 가중에 민감",
       "실측 사례: wikitext 단독 → 다운스트림 실패 (잘못된 M) / 장문 CoT 갭 → longcot 혼합 추가로 회복 (81.6 lossless)",
       "귀결: 데이터 레시피 설계 = M_eff의 경계 재가중 — 시행착오가 원리로 통역됨"], tfs=15)
# 9 claims ladder
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(.45), Inches(.2), Inches(12.4), Inches(.8))
p = tb.text_frame.paragraphs[0]; p.text = "8. 주장 사다리와 증명 상태"
p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = RGBColor(0x1a, 0x36, 0x5d)
body = s.shapes.add_textbox(Inches(.6), Inches(1.15), Inches(12.2), Inches(5.9))
tf = body.text_frame; tf.word_wrap = True
rows = [
 ("회전 불변성 (fresh=raw)", "정확한 정리 + 9B 실증 (정답 수 동일)", "확정"),
 ("prefix-loss → 고유기저 정렬", "선형 정리(Rippel'14) + 비선형은 실측 (탄력성 단조curves)", "강함"),
 ("스펙트럼 집중의 존재", "표현 비등방성의 상속(문헌) + GHOST phantom(외부) + 우리 33→81", "강함"),
 ("중요도 = loss-정의 (E-T1)", "명시 M(lag-0·시차)·열노름·질의에너지 전부 학습 R과 무상관(cos 0.44) — calibration류 실패의 뿌리", "확정(실측)"),
 ("일반화 (혼합→다운스트림)", "부분공간 가설류(작음) + E-T2 이식 실측(3도메인 tax 곡선 일치) + LM-only R 전이; D-K 갭 논증은 기각(gap@32≈0)", "확정(실측)"),
 ("메뉴 격자 (E-T2)", "off-menu 폭(8, 96)은 3도메인 공통 유료 — 학습물 = 전순서가 아닌 메뉴 지점의 flag", "관측"),
 ("경계 민감성", "wikitext 실패 / longcot 회복 — 재현 가능한 레시피 실측", "확정"),
 ("남은 이론 실험", "① ppl-acc 괴리 정량화(ghost k32: ppl 1.3×인데 GSM8K −26.5) ② M_task 정렬도(부분공간 각도) ③ 경계갭 없는 안정성의 대안 정리(flag 최적성)", "TODO"),
]
for i, (k, v, st) in enumerate(rows):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"[{st}]  {k}  —  {v}"
    p.font.size = Pt(15); p.space_after = Pt(12)
# 10 references
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(.45), Inches(.2), Inches(12.4), Inches(.8))
p = tb.text_frame.paragraphs[0]; p.text = "9. Background Reference Map (기둥별)"
p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = RGBColor(0x1a, 0x36, 0x5d)
body = s.shapes.add_textbox(Inches(.6), Inches(1.1), Inches(12.2), Inches(6.1))
tf = body.text_frame; tf.word_wrap = True
refs = [
 ("정렬 원리", "Rippel+ '14 Nested Dropout(=PCA) · Kusupati+ '22 MRL · Eckart–Young(저랭크 최적성)"),
 ("암묵적 고유분해", "Oja '82 (streaming PCA) · stochastic power iteration 문헌"),
 ("표현의 스펙트럼 집중", "Ethayarajh '19 (anisotropy) · Gao+ '19 representation degeneration"),
 ("고유공간 안정성(일반화)", "Davis–Kahan '70 sin-θ · Yu+ '15 (variant) · 행렬 농도부등식(Tropp '12)"),
 ("직교 파라미터화", "OFT (Qiu+ '23) · QR retraction (Stiefel manifold 최적화, Absil+ '08)"),
 ("상태-차원 축 선행", "SSE '25 (row-sparse, from-scratch) · GHOST '26 (정적 절단) · MoM '25 · MatMamba '24 · MatryoshkaKV '24"),
 ("query-aware 읽기(대조)", "Quest '24 (KV 페이지 선택 — 토큰축; 우리 실측: 상태축은 밀집이라 부적용)"),
 ("모델 축소의 고전", "Moore '81 balanced truncation (GHOST의 뿌리; 우리의 '지우지 말라'와 대비)"),
 ("아키텍처", "Dao&Gu '24 Mamba-2/SSD · GLA/GDN/KDA/GDN-2 · Nemotron-H"),
]
for i, (k, v) in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{k}:  {v}"
    p.font.size = Pt(14); p.space_after = Pt(10)

prs.save("../docs/theory_deck.pptx")
print("THEORY DECK DONE -> docs/theory_deck.pptx")
